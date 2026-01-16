"""
PPT解析器模块
负责解析上传的PPTX文件,提取幻灯片内容、版式、样式等信息
"""
import uuid
from typing import List, Dict, Any
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE_TYPE
from loguru import logger

from app.models.schemas import (
    PPTParseResult,
    SlideData,
    SlideType,
    LayoutInfo,
    StyleInfo
)


class PPTParser:
    """PPT解析器类"""

    def __init__(self):
        """初始化解析器"""
        self.supported_extensions = ['.pptx']

    def parse(self, file_path: str, ppt_id: str = None) -> PPTParseResult:
        """
        解析PPT文件

        Args:
            file_path: PPT文件路径
            ppt_id: PPT唯一标识,不提供则自动生成

        Returns:
            PPTParseResult: 解析结果
        """
        try:
            # 生成PPT ID
            if ppt_id is None:
                ppt_id = f"ppt_{uuid.uuid4().hex[:12]}"

            logger.info(f"开始解析PPT: {file_path}, ID: {ppt_id}")

            # 检查文件是否存在
            file_path_obj = Path(file_path)
            if not file_path_obj.exists():
                raise FileNotFoundError(f"PPT文件不存在: {file_path}")

            # 检查文件扩展名
            if file_path_obj.suffix.lower() not in self.supported_extensions:
                raise ValueError(f"不支持的文件格式: {file_path_obj.suffix}")

            # 打开PPT文件
            prs = Presentation(file_path)

            # 提取主题信息
            theme = self._extract_theme(prs)

            # 提取元数据
            metadata = self._extract_metadata(prs)

            # 解析每一页幻灯片
            slides = []
            for idx, slide in enumerate(prs.slides):
                slide_data = self._parse_slide(slide, idx)
                slides.append(slide_data)
                logger.debug(f"已解析第 {idx + 1}/{len(prs.slides)} 页")

            # 构建解析结果
            result = PPTParseResult(
                ppt_id=ppt_id,
                filename=file_path_obj.name,
                total_slides=len(slides),
                slides=slides,
                theme=theme,
                metadata=metadata
            )

            logger.info(f"PPT解析完成: {ppt_id}, 共 {len(slides)} 页")
            return result

        except Exception as e:
            logger.error(f"解析PPT失败: {file_path}, 错误: {str(e)}")
            raise

    def _parse_slide(self, slide, slide_index: int) -> SlideData:
        """
        解析单页幻灯片

        Args:
            slide: python-pptx的Slide对象
            slide_index: 幻灯片索引

        Returns:
            SlideData: 幻灯片数据
        """
        # 提取文本内容
        content = self._extract_text(slide)

        # 判断幻灯片类型
        slide_type = self._detect_slide_type(slide)

        # 提取版式信息
        layout_info = self._extract_layout_info(slide)

        # 提取样式信息
        style_info = self._extract_style_info(slide)

        # 提取图片信息
        images = self._extract_images(slide)

        # 提取图表信息
        charts = self._extract_charts(slide)

        return SlideData(
            slide_id=slide_index + 1,
            slide_index=slide_index,
            slide_type=slide_type,
            content=content,
            layout_info=layout_info,
            style_info=style_info,
            images=images,
            charts=charts
        )

    def _extract_text(self, slide) -> str:
        """提取幻灯片中的所有文本"""
        text_parts = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                text_parts.append(shape.text.strip())
        return "\n".join(text_parts)

    def _detect_slide_type(self, slide) -> SlideType:
        """
        检测幻灯片类型

        根据幻灯片的版式名称和内容判断类型
        """
        layout_name = slide.slide_layout.name.lower()

        if "title" in layout_name and "content" not in layout_name:
            return SlideType.TITLE

        # 检查是否包含图表
        has_chart = any(
            shape.shape_type == MSO_SHAPE_TYPE.CHART
            for shape in slide.shapes
            if hasattr(shape, 'shape_type')
        )
        if has_chart:
            return SlideType.CHART

        # 检查是否主要是图片
        image_count = sum(
            1 for shape in slide.shapes
            if hasattr(shape, 'shape_type') and shape.shape_type == MSO_SHAPE_TYPE.PICTURE
        )
        if image_count >= 2:
            return SlideType.IMAGE

        # 检查是否是混合内容
        has_text = any(hasattr(shape, "text") and shape.text for shape in slide.shapes)
        has_image = image_count > 0
        if has_text and (has_image or has_chart):
            return SlideType.MIXED

        return SlideType.CONTENT

    def _extract_layout_info(self, slide) -> LayoutInfo:
        """提取版式信息"""
        layout = slide.slide_layout
        placeholders = []

        for shape in slide.shapes:
            if shape.is_placeholder:
                placeholder_info = {
                    "type": shape.placeholder_format.type,
                    "name": shape.name,
                    "left": shape.left,
                    "top": shape.top,
                    "width": shape.width,
                    "height": shape.height
                }
                placeholders.append(placeholder_info)

        return LayoutInfo(
            layout_type=layout.name,
            width=slide.width,
            height=slide.height,
            placeholders=placeholders
        )

    def _extract_style_info(self, slide) -> StyleInfo:
        """提取样式信息"""
        style_info = StyleInfo()

        # 尝试从第一个文本框提取字体信息
        for shape in slide.shapes:
            if hasattr(shape, "text_frame") and shape.text_frame.text:
                text_frame = shape.text_frame
                if text_frame.paragraphs:
                    paragraph = text_frame.paragraphs[0]
                    if paragraph.runs:
                        run = paragraph.runs[0]
                        font = run.font

                        if font.name:
                            style_info.font_name = font.name
                        if font.size:
                            style_info.font_size = font.size.pt
                        if font.color and hasattr(font.color, 'rgb'):
                            try:
                                rgb = font.color.rgb
                                style_info.font_color = f"#{rgb[0]:02x}{rgb[1]:02x}{rgb[2]:02x}"
                            except:
                                pass

                        # 提取对齐方式
                        if paragraph.alignment:
                            style_info.alignment = str(paragraph.alignment)

                        break

        # 提取背景色
        if hasattr(slide, 'background') and slide.background:
            try:
                fill = slide.background.fill
                if fill.type == 1:  # SOLID
                    rgb = fill.fore_color.rgb
                    style_info.background_color = f"#{rgb[0]:02x}{rgb[1]:02x}{rgb[2]:02x}"
            except:
                pass

        return style_info

    def _extract_images(self, slide) -> List[Dict[str, Any]]:
        """提取图片信息"""
        images = []
        for shape in slide.shapes:
            if hasattr(shape, 'shape_type') and shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                image_info = {
                    "name": shape.name,
                    "left": shape.left,
                    "top": shape.top,
                    "width": shape.width,
                    "height": shape.height
                }
                images.append(image_info)
        return images

    def _extract_charts(self, slide) -> List[Dict[str, Any]]:
        """提取图表信息"""
        charts = []
        for shape in slide.shapes:
            if hasattr(shape, 'shape_type') and shape.shape_type == MSO_SHAPE_TYPE.CHART:
                chart_info = {
                    "name": shape.name,
                    "chart_type": str(shape.chart.chart_type) if hasattr(shape, 'chart') else "unknown",
                    "left": shape.left,
                    "top": shape.top,
                    "width": shape.width,
                    "height": shape.height
                }
                charts.append(chart_info)
        return charts

    def _extract_theme(self, prs: Presentation) -> Dict[str, Any]:
        """提取主题信息"""
        theme = {
            "slide_width": prs.slide_width,
            "slide_height": prs.slide_height,
        }

        # 尝试提取配色方案
        try:
            if hasattr(prs, 'slide_master') and prs.slide_master:
                theme["master_name"] = prs.slide_master.name
        except:
            pass

        return theme

    def _extract_metadata(self, prs: Presentation) -> Dict[str, Any]:
        """提取PPT元数据"""
        metadata = {}

        try:
            core_props = prs.core_properties
            if core_props.title:
                metadata["title"] = core_props.title
            if core_props.author:
                metadata["author"] = core_props.author
            if core_props.subject:
                metadata["subject"] = core_props.subject
            if core_props.created:
                metadata["created"] = str(core_props.created)
            if core_props.modified:
                metadata["modified"] = str(core_props.modified)
        except:
            pass

        return metadata
