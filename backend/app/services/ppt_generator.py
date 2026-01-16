"""
PPT自动生成器
基于优化方案自动生成或修改PPT文件
"""
from pathlib import Path
from typing import Optional
from pptx import Presentation
from pptx.util import Pt, Inches
from pptx.enum.text import PP_ALIGN
from loguru import logger

from app.models.schemas import (
    FinalOptimizationPlan,
    PPTParseResult,
    PPTGenerateResult,
    OptimizationDimension
)


class PPTGenerator:
    """PPT生成器类"""

    def __init__(self, output_dir: str = "./outputs"):
        """
        初始化生成器

        Args:
            output_dir: 输出目录
        """
        self.output_dir = Path(output_dir)
        self.output_dir.mkdir(parents=True, exist_ok=True)
        logger.info(f"PPT生成器初始化,输出目录: {self.output_dir}")

    async def generate(
        self,
        original_ppt_path: str,
        ppt_data: PPTParseResult,
        optimization_plan: FinalOptimizationPlan
    ) -> PPTGenerateResult:
        """
        生成优化后的PPT

        Args:
            original_ppt_path: 原始PPT路径
            ppt_data: PPT解析数据
            optimization_plan: 优化方案

        Returns:
            PPTGenerateResult: 生成结果
        """
        try:
            logger.info(f"开始生成优化PPT: {optimization_plan.ppt_id}")

            # 优先尝试使用python-pptx库生成
            # 在生产环境中,这里应该优先调用通义千问PPT生成API
            result = await self._generate_with_library(
                original_ppt_path,
                ppt_data,
                optimization_plan
            )

            logger.info(f"PPT生成完成: {result.output_filename}")
            return result

        except Exception as e:
            logger.error(f"PPT生成失败: {str(e)}")
            return PPTGenerateResult(
                ppt_id=optimization_plan.ppt_id,
                output_filename="",
                output_path="",
                generation_method="failed",
                success=False,
                error_message=str(e)
            )

    async def _generate_with_library(
        self,
        original_path: str,
        ppt_data: PPTParseResult,
        plan: FinalOptimizationPlan
    ) -> PPTGenerateResult:
        """
        使用python-pptx库生成PPT

        Args:
            original_path: 原始PPT路径
            ppt_data: PPT数据
            plan: 优化方案

        Returns:
            PPTGenerateResult: 生成结果
        """
        # 打开原始PPT
        prs = Presentation(original_path)

        # 应用优化建议
        self._apply_optimizations(prs, plan)

        # 保存新PPT
        output_filename = f"optimized_{ppt_data.filename}"
        output_path = self.output_dir / output_filename
        prs.save(str(output_path))

        return PPTGenerateResult(
            ppt_id=plan.ppt_id,
            output_filename=output_filename,
            output_path=str(output_path),
            generation_method="library",
            success=True
        )

    def _apply_optimizations(self, prs: Presentation, plan: FinalOptimizationPlan):
        """
        应用优化建议到PPT

        Args:
            prs: Presentation对象
            plan: 优化方案
        """
        # 按页面分组建议
        suggestions_by_slide = {}
        for sug in plan.suggestions:
            if sug.slide_index not in suggestions_by_slide:
                suggestions_by_slide[sug.slide_index] = []
            suggestions_by_slide[sug.slide_index].append(sug)

        # 遍历每一页应用优化
        for slide_idx, suggestions in suggestions_by_slide.items():
            if slide_idx >= len(prs.slides):
                continue

            slide = prs.slides[slide_idx]
            for sug in suggestions:
                self._apply_suggestion(slide, sug)

    def _apply_suggestion(self, slide, suggestion):
        """
        应用单个优化建议

        Args:
            slide: Slide对象
            suggestion: 优化建议
        """
        try:
            dimension = suggestion.optimization_dimension

            if dimension == OptimizationDimension.FONT:
                self._apply_font_optimization(slide, suggestion)
            elif dimension == OptimizationDimension.COLOR:
                self._apply_color_optimization(slide, suggestion)
            elif dimension == OptimizationDimension.LAYOUT:
                self._apply_layout_optimization(slide, suggestion)
            elif dimension == OptimizationDimension.CONTENT:
                self._apply_content_optimization(slide, suggestion)

        except Exception as e:
            logger.warning(f"应用建议失败: {str(e)}")

    def _apply_font_optimization(self, slide, suggestion):
        """应用字体优化"""
        # 简化实现: 修改所有文本框的字体
        for shape in slide.shapes:
            if hasattr(shape, "text_frame"):
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        # 这里应该根据suggestion的具体内容解析字体设置
                        # 简化处理: 设置默认字体大小
                        if run.font:
                            run.font.size = Pt(18)

    def _apply_color_optimization(self, slide, suggestion):
        """应用配色优化"""
        # 简化实现
        pass

    def _apply_layout_optimization(self, slide, suggestion):
        """应用版式优化"""
        # 简化实现
        pass

    def _apply_content_optimization(self, slide, suggestion):
        """应用内容优化"""
        # 简化实现: 可以修改文本内容
        pass
